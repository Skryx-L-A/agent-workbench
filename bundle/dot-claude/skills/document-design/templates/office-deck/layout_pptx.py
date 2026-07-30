"""layout_pptx.py -- the MECHANICS of a PPTX deck: slide size, blank-canvas text
boxes with explicit fonts (never the inherited PowerPoint theme), a hard character
budget checked after the content is assembled, font embedding. Mirrors
templates/deck/theme.typ one level down: same split (mechanics here, decisions in
tokens.py, content in deck.py), same gate.

Built on blank slide layouts rather than the default template's placeholders: the
default template's placeholders inherit Calibri/Office-theme fonts silently unless
every run is overridden anyway, and a corporate .potx (a client's own template) is a
different situation entirely -- see SKILL.md "Corporate templates" for that path,
which fills existing placeholders instead of building shapes.
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

_here = Path(__file__).resolve().parent
for _candidate in (_here / "scripts", _here.parent / "scripts", _here.parent.parent / "scripts"):
    if (_candidate / "officefonts.py").exists():
        sys.path.insert(0, str(_candidate))
        break
else:
    sys.exit(
        "officefonts.py not found in ./scripts, ../scripts or ../../scripts next to this "
        "template.\nCopy scripts/officefonts.py alongside this template, or run it from "
        "inside the document-design project."
    )
import officefonts  # noqa: E402

GATE_MESSAGE = (
    "Layout for this deck has not been decided yet.\n"
    "Open tokens.py, replace every value with your own decision for THIS deck,\n"
    "write the reason on each `# why:` comment, then set DESIGN_DECIDED = True.\n"
    "Do not simply flip the flag: the shipped numbers are an example, not a house style."
)


def gate(tokens) -> None:
    if not getattr(tokens, "DESIGN_DECIDED", False):
        sys.exit(GATE_MESSAGE)
    for group in (tokens.FONT_TEXT, tokens.FONT_DISPLAY):
        if group["family"] == "CHOOSE-ME" or not Path(group["regular"]).exists():
            sys.exit(
                f"Font not decided or not found: {group}.\n"
                "Run scripts/make-static-fonts.sh on a variable font first, or point\n"
                "FONT_TEXT/FONT_DISPLAY at real Regular/Bold .ttf files."
            )


def _rgb(hex_str: str) -> RGBColor:
    return RGBColor.from_string(hex_str)


class BudgetError(SystemExit):
    pass


class TITLE:
    def __init__(self, title: str, subtitle: str | None = None, notes: str | None = None):
        self.title, self.subtitle, self.notes = title, subtitle, notes


class SECTION:
    def __init__(self, title: str, notes: str | None = None):
        self.title, self.notes = title, notes


class CONTENT:
    def __init__(self, heading: str, bullets: list[str], notes: str | None = None):
        self.heading, self.bullets, self.notes = heading, bullets, notes


class IMAGE:
    def __init__(self, heading: str, path: str, caption: str | None = None, notes: str | None = None):
        self.heading, self.path, self.caption, self.notes = heading, path, caption, notes


def _check_budgets(tokens, slides: list) -> list[str]:
    """Returns violation strings instead of raising per-slide -- one loud report at
    the end beats stopping at the first slide and hiding the rest."""
    b = tokens.BUDGET_CHARS
    problems = []
    for i, s in enumerate(slides, start=1):
        heading = getattr(s, "heading", None) or getattr(s, "title", None)
        if heading and len(heading) > b["heading"]:
            problems.append(f"slide {i}: heading is {len(heading)} chars, budget is {b['heading']}: {heading!r}")
        bullets = getattr(s, "bullets", None)
        if bullets:
            if len(bullets) > b["bullets_per_slide"]:
                problems.append(
                    f"slide {i}: {len(bullets)} bullets, budget is {b['bullets_per_slide']} -- split into two slides"
                )
            for j, bullet in enumerate(bullets, start=1):
                if len(bullet) > b["bullet"]:
                    problems.append(f"slide {i} bullet {j}: {len(bullet)} chars, budget is {b['bullet']}: {bullet!r}")
        if s.notes and len(s.notes) > b["notes"]:
            problems.append(f"slide {i}: notes are {len(s.notes)} chars, budget is {b['notes']}")
    return problems


def _textbox(slide, tokens, left_in, top_in, width_in, height_in, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in))
    tf = box.text_frame
    tf.word_wrap = True
    from pptx.enum.text import MSO_AUTO_SIZE

    tf.auto_size = MSO_AUTO_SIZE.NONE  # overflow must be visible, never silently shrunk
    tf.vertical_anchor = anchor
    for p in tf.paragraphs:
        p.alignment = align
    return box, tf


def _run(paragraph, text, tokens, family, size_pt, color, bold=False, italic=False):
    run = paragraph.add_run()
    run.text = text
    run.font.name = family
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = _rgb(color)
    return run


def _blank_slide(prs):
    # Layout 6 is "Blank" in python-pptx's default template -- no inherited
    # placeholder, no inherited theme font silently applied to unset runs.
    return prs.slides.add_slide(prs.slide_layouts[6])


def _build_title(prs, tokens, block: TITLE):
    slide = _blank_slide(prs)
    m = tokens.MARGIN_MM
    left_in = m["left"] / 25.4
    width_in = tokens.SLIDE_WIDTH_IN - 2 * left_in
    _, tf = _textbox(slide, tokens, left_in, tokens.SLIDE_HEIGHT_IN / 2 - 0.9, width_in, 1.2, anchor=MSO_ANCHOR.BOTTOM)
    _run(tf.paragraphs[0], block.title, tokens, tokens.FONT_DISPLAY["family"], tokens.SIZE_PT["title"], tokens.INK["text"], bold=True)
    if block.subtitle:
        _, tf2 = _textbox(slide, tokens, left_in, tokens.SLIDE_HEIGHT_IN / 2 + 0.35, width_in, 0.6)
        _run(tf2.paragraphs[0], block.subtitle, tokens, tokens.FONT_TEXT["family"], tokens.SIZE_PT["body"], tokens.INK["soft"], italic=True)
    return slide


def _build_section(prs, tokens, block: SECTION):
    slide = _blank_slide(prs)
    m = tokens.MARGIN_MM
    left_in = m["left"] / 25.4
    width_in = tokens.SLIDE_WIDTH_IN - 2 * left_in
    # A section opener carries one visual gesture (the accent rule), not decoration.
    # A 1-cell table -- the trick that draws the report template's title-page rule --
    # does NOT work here: PowerPoint tables enforce a minimum row height from the
    # cell's default font size, so a Pt(3) row request rendered as a ~0.3in block that
    # sat on top of the title text (found rendering this template's own section
    # slide). A plain rectangle autoshape has no such floor.
    from pptx.enum.shapes import MSO_SHAPE

    rule = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left_in), Inches(tokens.SLIDE_HEIGHT_IN / 2 - 0.55), Inches(1.2), Pt(6)
    )
    rule.fill.solid()
    rule.fill.fore_color.rgb = _rgb(tokens.INK["accent"])
    rule.line.fill.background()
    rule.shadow.inherit = False

    _, tf = _textbox(slide, tokens, left_in, tokens.SLIDE_HEIGHT_IN / 2 - 0.35, width_in, 1.0)
    _run(tf.paragraphs[0], block.title, tokens, tokens.FONT_DISPLAY["family"], tokens.SIZE_PT["title"], tokens.INK["text"], bold=True)
    return slide


def _build_content(prs, tokens, block: CONTENT):
    slide = _blank_slide(prs)
    m = tokens.MARGIN_MM
    left_in = m["left"] / 25.4
    top_in = m["top"] / 25.4
    width_in = tokens.SLIDE_WIDTH_IN - 2 * left_in

    _, tf = _textbox(slide, tokens, left_in, top_in, width_in, 0.9)
    _run(tf.paragraphs[0], block.heading, tokens, tokens.FONT_DISPLAY["family"], tokens.SIZE_PT["heading"], tokens.INK["text"], bold=True)

    body_top = top_in + 1.0
    body_height = tokens.SLIDE_HEIGHT_IN - body_top - (m["bottom"] / 25.4)
    _, body_tf = _textbox(slide, tokens, left_in, body_top, width_in, body_height)
    for i, bullet in enumerate(block.bullets):
        p = body_tf.paragraphs[0] if i == 0 else body_tf.add_paragraph()
        p.space_after = Pt(14)
        _run(p, f"–  {bullet}", tokens, tokens.FONT_TEXT["family"], tokens.SIZE_PT["body"], tokens.INK["text"])
    return slide


def _build_image(prs, tokens, block: IMAGE):
    slide = _blank_slide(prs)
    m = tokens.MARGIN_MM
    left_in = m["left"] / 25.4
    top_in = m["top"] / 25.4
    width_in = tokens.SLIDE_WIDTH_IN - 2 * left_in

    _, tf = _textbox(slide, tokens, left_in, top_in, width_in, 0.7)
    _run(tf.paragraphs[0], block.heading, tokens, tokens.FONT_DISPLAY["family"], tokens.SIZE_PT["heading"], tokens.INK["text"], bold=True)

    img_top = top_in + 0.85
    img_height = tokens.SLIDE_HEIGHT_IN - img_top - (m["bottom"] / 25.4) - (0.35 if block.caption else 0)
    slide.shapes.add_picture(block.path, Inches(left_in), Inches(img_top), height=Inches(img_height))

    if block.caption:
        _, cap_tf = _textbox(slide, tokens, left_in, tokens.SLIDE_HEIGHT_IN - (m["bottom"] / 25.4) - 0.3, width_in, 0.3)
        _run(cap_tf.paragraphs[0], block.caption, tokens, tokens.FONT_TEXT["family"], tokens.SIZE_PT["small"], tokens.INK["soft"])
    return slide


_DISPATCH = {TITLE: _build_title, SECTION: _build_section, CONTENT: _build_content, IMAGE: _build_image}


def build(tokens, slides: list, out_path: str):
    gate(tokens)

    problems = _check_budgets(tokens, slides)
    if problems:
        sys.exit(
            "Character budget exceeded (tokens.BUDGET_CHARS) -- PPTX text does not "
            "reflow, an over-budget placeholder overflows the slide invisibly in the "
            "source:\n  " + "\n  ".join(problems)
        )

    prs = Presentation()
    prs.slide_width = Inches(tokens.SLIDE_WIDTH_IN)
    prs.slide_height = Inches(tokens.SLIDE_HEIGHT_IN)

    for block in slides:
        slide = _DISPATCH[type(block)](prs, tokens, block)
        if block.notes:
            slide.notes_slide.notes_text_frame.text = block.notes

    Path(out_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)

    font_specs = []
    for group in (tokens.FONT_TEXT, tokens.FONT_DISPLAY):
        spec = {"family": group["family"], "regular": group["regular"]}
        if group.get("bold"):
            spec["bold"] = group["bold"]
        font_specs.append(spec)
    seen = {}
    for spec in font_specs:
        seen[spec["family"]] = spec
    officefonts.embed_fonts_pptx(out_path, list(seen.values()))
